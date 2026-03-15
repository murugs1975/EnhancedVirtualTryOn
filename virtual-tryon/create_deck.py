from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand palette ──────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x3E)   # dark background
ACCENT     = RGBColor(0x00, 0xC2, 0xFF)   # cyan / 6G feel
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE8, 0xEC, 0xF0)
MID_GRAY   = RGBColor(0x8A, 0x99, 0xAA)
GOLD       = RGBColor(0xFF, 0xC1, 0x07)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ══════════════════════════════════════════════════════════════
# Helpers
# ══════════════════════════════════════════════════════════════

def add_rect(slide, l, t, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    return txBox


def add_bullet_block(slide, items, l, t, w, h,
                     font_size=16, color=WHITE,
                     bullet_color=None, indent=Inches(0.25)):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = True
    bc = bullet_color or ACCENT

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # bullet character
        brun = p.add_run()
        brun.text = "▸  "
        brun.font.size  = Pt(font_size)
        brun.font.color.rgb = bc
        brun.font.bold  = True
        # content
        run = p.add_run()
        run.text = item
        run.font.size  = Pt(font_size)
        run.font.color.rgb = color
    return txBox


def add_divider(slide, t, color=ACCENT, thickness=3):
    line = slide.shapes.add_shape(1, Inches(0.5), t, Inches(12.33), Pt(thickness))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Hero
# ══════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)

# Full dark background
add_rect(s1, 0, 0, SLIDE_W, SLIDE_H, NAVY)

# Accent bar — left edge
add_rect(s1, 0, 0, Inches(0.18), SLIDE_H, ACCENT)

# Top accent strip
add_rect(s1, Inches(0.18), 0, SLIDE_W, Inches(0.08), ACCENT)

# Course chip
chip = add_rect(s1, Inches(0.5), Inches(0.3), Inches(3.2), Inches(0.45), ACCENT)
add_text(s1, "MSDS · COMPUTERVISION 462  |  Final Project",
         Inches(0.55), Inches(0.3), Inches(3.1), Inches(0.45),
         font_size=9, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

# Main title
add_text(s1,
         "Enhanced Virtual Try-On\nin the Metaverse",
         Inches(0.5), Inches(1.1), Inches(8), Inches(2.0),
         font_size=48, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Subtitle
add_text(s1,
         "Leveraging 6G Technology for High-Resolution\nCloth Detection & Fitting Room Experiences",
         Inches(0.5), Inches(3.0), Inches(8.5), Inches(1.2),
         font_size=20, bold=False, color=ACCENT, align=PP_ALIGN.LEFT)

# Thin separator
add_rect(s1, Inches(0.5), Inches(4.3), Inches(5), Inches(0.04), MID_GRAY)

# Team & date line
add_text(s1, "Team: Joyati  ·  Biraj  ·  Murughanandam S.",
         Inches(0.5), Inches(4.45), Inches(7), Inches(0.4),
         font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)
add_text(s1, "February 2026",
         Inches(0.5), Inches(4.85), Inches(3), Inches(0.4),
         font_size=12, color=MID_GRAY, align=PP_ALIGN.LEFT)

# Right-side decorative "6G" watermark
add_text(s1, "6G", Inches(9.8), Inches(1.5), Inches(3), Inches(3),
         font_size=160, bold=True,
         color=RGBColor(0x1A, 0x2D, 0x55),   # very dark, almost invisible
         align=PP_ALIGN.CENTER)
add_text(s1, "⬡", Inches(10.3), Inches(3.6), Inches(2), Inches(1.5),
         font_size=80, bold=False,
         color=RGBColor(0x00, 0x4A, 0x6E),
         align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ══════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
add_rect(s2, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s2, 0, 0, Inches(0.18), SLIDE_H, ACCENT)

# Header band
add_rect(s2, Inches(0.18), 0, SLIDE_W, Inches(1.1), RGBColor(0x0A, 0x14, 0x30))
add_text(s2, "01", Inches(0.3), Inches(0.1), Inches(0.9), Inches(0.9),
         font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
add_text(s2, "Problem Statement",
         Inches(1.3), Inches(0.25), Inches(8), Inches(0.65),
         font_size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

add_divider(s2, Inches(1.12))

# Two-column layout
# LEFT column — challenge context
add_text(s2, "CONTEXT",
         Inches(0.5), Inches(1.25), Inches(5.5), Inches(0.35),
         font_size=11, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)

add_text(s2,
         "Online fashion retail demands immersive, realistic try-on experiences — "
         "yet current virtual fitting rooms fall short on resolution, realism, and scale.",
         Inches(0.5), Inches(1.6), Inches(5.6), Inches(1.0),
         font_size=15, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

challenges = [
    "Low synthesized image resolution (256×192) creates noticeable artifacts",
    "Misalignment between warped clothing and desired body regions",
    "Existing architectures struggle with texture sharpness at high res",
    "Scalability bottleneck as Metaverse & 6G demand real-time, 4K-class fidelity",
]
add_bullet_block(s2, challenges,
                 Inches(0.5), Inches(2.65), Inches(5.6), Inches(3.0),
                 font_size=15)

# RIGHT column — our focus
add_rect(s2, Inches(6.6), Inches(1.2), Inches(6.2), Inches(5.8),
         RGBColor(0x0A, 0x14, 0x30))

add_text(s2, "OUR FOCUS",
         Inches(6.85), Inches(1.35), Inches(5.7), Inches(0.35),
         font_size=11, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

focus_points = [
    ("High-Res Try-On",
     "Synthesize 1024×768 (and beyond) virtual try-on images using VITON-HD-class techniques."),
    ("Cloth Detection",
     "Robust segmentation & detection pipeline for accurate clothing region identification."),
    ("6G + Metaverse",
     "Exploit ultra-low latency & high bandwidth of 6G to deliver real-time immersive fitting rooms."),
]

y = Inches(1.85)
for title, desc in focus_points:
    add_rect(s2, Inches(6.75), y, Inches(0.05), Inches(0.6), ACCENT)
    add_text(s2, title, Inches(6.95), y, Inches(5.3), Inches(0.35),
             font_size=14, bold=True, color=ACCENT)
    add_text(s2, desc, Inches(6.95), y + Inches(0.32), Inches(5.3), Inches(0.55),
             font_size=13, color=LIGHT_GRAY)
    y += Inches(1.3)


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — Team Profile
# ══════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
add_rect(s3, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s3, 0, 0, Inches(0.18), SLIDE_H, ACCENT)
add_rect(s3, Inches(0.18), 0, SLIDE_W, Inches(1.1), RGBColor(0x0A, 0x14, 0x30))

add_text(s3, "02", Inches(0.3), Inches(0.1), Inches(0.9), Inches(0.9),
         font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
add_text(s3, "Team Profile",
         Inches(1.3), Inches(0.25), Inches(8), Inches(0.65),
         font_size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_divider(s3, Inches(1.12))

team = [
    {
        "name": "Joyati",
        "role": "ML Engineer",
        "focus": "Model Testing & Monitoring",
        "details": [
            "Spearheads model testing & monitoring strategies",
            "Builds evaluation pipelines in Python",
            "Defines QA metrics for try-on synthesis quality",
        ],
        "icon": "🔬",
    },
    {
        "name": "Biraj Mishra",
        "role": "Software Engineer",
        "focus": "Systems & Integration",
        "details": [
            "Software engineering background",
            "System architecture & API design",
            "Integrates CV models into production pipelines",
        ],
        "icon": "⚙️",
    },
    {
        "name": "Murughanandam S.",
        "role": "Senior Director",
        "focus": "Research & Strategy",
        "details": [
            "Senior industry leadership perspective",
            "Guides research direction & benchmarking",
            "Connects academic work to real-world deployment",
        ],
        "icon": "🎯",
    },
]

card_w = Inches(3.7)
card_h = Inches(5.2)
card_t = Inches(1.3)
gaps   = [Inches(0.45), Inches(4.35), Inches(8.25)]

for i, (member, left) in enumerate(zip(team, gaps)):
    # Card background
    add_rect(s3, left, card_t, card_w, card_h, RGBColor(0x0A, 0x14, 0x30))
    # Top accent strip on each card
    add_rect(s3, left, card_t, card_w, Inches(0.07), ACCENT)

    # Avatar circle placeholder (just a colored square here)
    add_rect(s3, left + Inches(1.35), card_t + Inches(0.2),
             Inches(1.0), Inches(1.0), RGBColor(0x00, 0x4A, 0x6E))
    add_text(s3, member["icon"],
             left + Inches(1.35), card_t + Inches(0.2),
             Inches(1.0), Inches(1.0),
             font_size=32, align=PP_ALIGN.CENTER)

    # Name
    add_text(s3, member["name"],
             left + Inches(0.1), card_t + Inches(1.35),
             card_w - Inches(0.2), Inches(0.45),
             font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Role chip
    add_rect(s3, left + Inches(0.6), card_t + Inches(1.82),
             card_w - Inches(1.2), Inches(0.32), ACCENT)
    add_text(s3, member["role"],
             left + Inches(0.6), card_t + Inches(1.82),
             card_w - Inches(1.2), Inches(0.32),
             font_size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Focus label
    add_text(s3, member["focus"].upper(),
             left + Inches(0.15), card_t + Inches(2.28),
             card_w - Inches(0.3), Inches(0.28),
             font_size=9, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Divider
    add_rect(s3, left + Inches(0.3), card_t + Inches(2.6),
             card_w - Inches(0.6), Inches(0.025), MID_GRAY)

    # Detail bullets
    add_bullet_block(s3, member["details"],
                     left + Inches(0.15), card_t + Inches(2.7),
                     card_w - Inches(0.3), Inches(2.2),
                     font_size=12, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — Dataset
# ══════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
add_rect(s4, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s4, 0, 0, Inches(0.18), SLIDE_H, ACCENT)
add_rect(s4, Inches(0.18), 0, SLIDE_W, Inches(1.1), RGBColor(0x0A, 0x14, 0x30))

add_text(s4, "03", Inches(0.3), Inches(0.1), Inches(0.9), Inches(0.9),
         font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
add_text(s4, "Dataset",
         Inches(1.3), Inches(0.25), Inches(8), Inches(0.65),
         font_size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_divider(s4, Inches(1.12))

# Dataset name badge
add_rect(s4, Inches(0.45), Inches(1.25), Inches(7.5), Inches(0.52),
         RGBColor(0x0A, 0x14, 0x30))
add_text(s4, "VITON-HD  ·  High-Resolution Zalando Dataset  (Kaggle)",
         Inches(0.55), Inches(1.27), Inches(7.3), Inches(0.48),
         font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)

# Description paragraph
add_text(s4,
         "Image-based virtual try-on dataset designed to transfer a target clothing item onto "
         "the corresponding region of a person. VITON-HD addresses the critical limitations of "
         "prior low-resolution (256×192) benchmarks by providing 1024×768 paired person–garment images.",
         Inches(0.45), Inches(1.85), Inches(7.5), Inches(1.1),
         font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

# Key stats row
stats = [
    ("1024 × 768", "Image Resolution"),
    ("13,679", "Training Pairs"),
    ("2,032", "Test Pairs"),
    ("Zalando", "Source Platform"),
]

stat_w = Inches(2.8)
stat_t = Inches(3.1)
for j, (val, lbl) in enumerate(stats):
    sx = Inches(0.45) + j * Inches(3.05)
    add_rect(s4, sx, stat_t, stat_w, Inches(1.05), RGBColor(0x0A, 0x14, 0x30))
    add_rect(s4, sx, stat_t, stat_w, Inches(0.05), ACCENT)
    add_text(s4, val, sx, stat_t + Inches(0.1), stat_w, Inches(0.55),
             font_size=26, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s4, lbl, sx, stat_t + Inches(0.65), stat_w, Inches(0.35),
             font_size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

# Why this dataset
add_text(s4, "WHY THIS DATASET",
         Inches(0.45), Inches(4.35), Inches(7.5), Inches(0.3),
         font_size=11, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

why_items = [
    "State-of-the-art resolution benchmark — enables rigorous high-res synthesis evaluation",
    "Includes segmentation maps & body keypoints essential for ALIAS normalization training",
    "Paired person + garment images enable supervised cloth warping & fusion experiments",
    "Publicly available on Kaggle — reproducible & community-validated splits",
]
add_bullet_block(s4, why_items,
                 Inches(0.45), Inches(4.65), Inches(7.5), Inches(2.5),
                 font_size=14, color=LIGHT_GRAY)

# Right panel — method highlight
add_rect(s4, Inches(8.3), Inches(1.2), Inches(4.7), Inches(6.0),
         RGBColor(0x0A, 0x14, 0x30))
add_text(s4, "VITON-HD PIPELINE",
         Inches(8.4), Inches(1.35), Inches(4.5), Inches(0.35),
         font_size=11, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)

pipeline_steps = [
    ("1", "Segmentation Map\nPreparation"),
    ("2", "Clothing Item\nCoarse Fitting"),
    ("3", "ALIAS Normalization\n(Misalignment Handling)"),
    ("4", "ALIAS Generator\n(1024×768 Output)"),
]
py = Inches(1.85)
for num, desc in pipeline_steps:
    add_rect(s4, Inches(8.4), py, Inches(0.45), Inches(0.45),
             ACCENT)
    add_text(s4, num, Inches(8.4), py, Inches(0.45), Inches(0.45),
             font_size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s4, desc, Inches(8.95), py, Inches(3.8), Inches(0.6),
             font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)
    if num != "4":
        add_rect(s4, Inches(8.6), py + Inches(0.48), Inches(0.05), Inches(0.18), MID_GRAY)
    py += Inches(1.05)

add_text(s4, "Surpasses all baselines in FID & SSIM",
         Inches(8.4), Inches(6.15), Inches(4.5), Inches(0.4),
         font_size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER, italic=True)

# ── Save ────────────────────────────────────────────────────────
out = "VirtualTryOn_Presentation.pptx"
prs.save(out)
print(f"Saved → {out}")
